import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Define paths
PPTX_PATH = 'website-01.pptx'
ASSETS_DIR = 'assets'
DATA_FILE = 'slides_data.json'

# Ensure assets directory exists
if not os.path.exists(ASSETS_DIR):
    os.makedirs(ASSETS_DIR)

def extract_text_from_frame(text_frame):
    """
    Extracts text and structural information (like paragraphs and bullet levels) from a text frame.
    """
    paragraphs_data = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        paragraphs_data.append({
            'text': text,
            'level': paragraph.level,  # indent level, useful for bullets
            'is_bullet': bool(paragraph.level > 0 or (paragraph.text and paragraph.text.startswith(('-', '•'))))
        })
    return paragraphs_data

def process_shapes(shapes, slide_num, slide_data):
    """
    Recursively processes shapes on a slide to extract text, tables, and images.
    """
    for i, shape in enumerate(shapes):
        shape_type = shape.shape_type
        
        # 1. Handle Group Shapes (Recurse)
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            process_shapes(shape.shapes, slide_num, slide_data)
            continue
            
        # 2. Handle Table Shapes
        if shape.has_table:
            table = shape.table
            table_data = []
            for r_idx, row in enumerate(table.rows):
                row_data = []
                for c_idx, cell in enumerate(row.cells):
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            slide_data['tables'].append(table_data)
            continue

        # 3. Handle Picture Shapes
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.ext.lower()
            
            # Generate a clean unique name
            img_filename = f"slide_{slide_num}_img_{len(slide_data['images']) + 1}.{ext}"
            img_path = os.path.join(ASSETS_DIR, img_filename)
            
            with open(img_path, 'wb') as f:
                f.write(image.blob)
                
            slide_data['images'].append({
                'filename': img_filename,
                'path': img_path,
                'original_name': image.filename
            })
            continue

        # 4. Handle Text Box / Auto Shapes with Text
        if shape.has_text_frame:
            paragraphs = extract_text_from_frame(shape.text_frame)
            if paragraphs:
                # Classify if this is likely a title or general content
                is_title = False
                try:
                    if shape == shape.part.slide.shapes.title:
                        is_title = True
                except:
                    pass
                
                # Check text length / capitalization to guess if it is a title
                first_txt = paragraphs[0]['text']
                if not is_title and len(first_txt) < 60 and first_txt.isupper():
                    is_title = True
                
                if is_title and not slide_data['title']:
                    slide_data['title'] = first_txt
                    # Add rest of paragraphs to text blocks
                    if len(paragraphs) > 1:
                        slide_data['text_blocks'].extend(paragraphs[1:])
                else:
                    slide_data['text_blocks'].extend(paragraphs)

def main():
    if not os.path.exists(PPTX_PATH):
        print(f"Error: {PPTX_PATH} not found.")
        return

    prs = Presentation(PPTX_PATH)
    all_slides = []

    print(f"Processing {len(prs.slides)} slides...")
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'text_blocks': [],
            'tables': [],
            'images': []
        }
        
        # Check standard slide title shape
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text.strip()
            
        # Process all shapes (recursively handling groups)
        process_shapes(slide.shapes, slide_num, slide_data)
        
        # Clean up title if it ended up in text blocks
        if slide_data['title']:
            slide_data['text_blocks'] = [
                tb for tb in slide_data['text_blocks'] 
                if tb['text'] != slide_data['title']
            ]
            
        print(f"Slide {slide_num}: Title='{slide_data['title']}', Text Blocks={len(slide_data['text_blocks'])}, Tables={len(slide_data['tables'])}, Images={len(slide_data['images'])}")
        all_slides.append(slide_data)
        
    # Save slide data to JSON
    with open(DATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(all_slides, f, indent=4, ensure_ascii=False)
        
    print(f"\nExtraction complete! Slide content saved to {DATA_FILE}")
    print(f"All images saved to {ASSETS_DIR}/")

if __name__ == '__main__':
    main()
